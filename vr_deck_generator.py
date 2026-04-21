"""
VR Deck Generator — Powerplay Customer Success
================================================
Reusable engine that turns an AA_Event depth-funnel CSV into a polished
Value Review deck for any org.

QUICK USAGE
-----------
    from vr_deck_generator import generate_vr_deck

    # Generate decks for EVERY org in the CSV (no customization)
    generate_vr_deck(csv_path='usage.csv', output_dir='./decks')

    # Generate one deck for one org, fully customized
    generate_vr_deck(
        csv_path='usage.csv',
        org_filter='G D CONSTRUCTION & CO',
        output_dir='./decks',
        customizations={
            'G D CONSTRUCTION & CO': {
                'owner_name': 'Mitranjan Ganguly',
                'owner_title': 'Owner',
                'onboarding_goal': 'We want to track labour attendance across our sites.',
                'onboarding_goal_short': 'labour attendance tracking',
                'renewal_ask': 'same_arr',
                'pain_points': ['Material leakage across sites',
                                'Dormant site visibility'],
                'kam_name': 'Altaf',
                'trainer_name': 'Wasim',
                'year1_comparison': {'events': 20000, 'users': 25, 'sites': 15},
            }
        }
    )

CSV REQUIREMENTS
----------------
Expected columns (from AA_Event depth funnel export):
  project_name, role, user_name, org_id, org_name, ARR, status,
  subs_start_date, renewal_month, module, event_type,
  event_date, event_name, event_count, activity_date,
  kam_email, trainer_email, min_start_date, max_end_date

AUTHOR: Powerplay CS Team
"""

import os
import re
import json
import zipfile
import datetime as dt
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION


# =============================================================================
# PALETTE
# =============================================================================
NAVY       = RGBColor(0x1E, 0x27, 0x61)
NAVY_DEEP  = RGBColor(0x15, 0x1C, 0x48)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_SOFT = RGBColor(0xFE, 0xF3, 0xC7)
CREAM      = RGBColor(0xF8, 0xFA, 0xFC)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
INK        = RGBColor(0x0F, 0x17, 0x2A)
MUTED      = RGBColor(0x64, 0x74, 0x8B)
BORDER     = RGBColor(0xE2, 0xE8, 0xF0)
SUCCESS    = RGBColor(0x10, 0xB9, 0x81)
WARN       = RGBColor(0xEF, 0x44, 0x44)
LIGHT_BLUE = RGBColor(0xCA, 0xDC, 0xFC)


# =============================================================================
# PPTX HELPERS
# =============================================================================
def _add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    return shp


def _add_text(slide, x, y, w, h, text, *, font="Calibri", size=14, bold=False,
              color=INK, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
              italic=False, spacing_after=0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = valign
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing_after:
            p.space_after = Pt(spacing_after)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def _add_circle(slide, cx, cy, d, fill, line=None):
    x = cx - d / 2
    y = cy - d / 2
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    return shp


def _fmt_int(n):
    try:
        return f"{int(n):,}"
    except Exception:
        return str(n)


def _safe_filename(name):
    """Turn an org name into a safe filename fragment."""
    s = re.sub(r"[^A-Za-z0-9]+", "_", name).strip("_")
    return s[:60] if len(s) > 60 else s


# =============================================================================
# STATS COMPILATION — reads CSV, computes stats for one org
# =============================================================================
def compile_stats(df, org_name, as_of_date=None):
    """
    Given a DataFrame filtered (or filterable) by org and an org name,
    return a stats dict matching the deck's expectations.
    """
    org_df = df[df['org_name'] == org_name].copy()
    if org_df.empty:
        raise ValueError(f"No rows found for org '{org_name}'")

    # Normalize event_count to int — sometimes comes as string or float
    org_df['event_count'] = pd.to_numeric(org_df['event_count'], errors='coerce').fillna(0)

    # Dates — try ISO first, fall back to dateutil if mixed formats
    org_df['activity_date'] = pd.to_datetime(
        org_df['activity_date'], errors='coerce', format='mixed'
    )
    # Cap analysis date: use caller's as_of_date if provided, else today.
    # This filters out any forecast/future-dated rows in the CSV (common with
    # depth-funnel exports that project dates beyond the run date).
    as_of = pd.to_datetime(as_of_date) if as_of_date else pd.Timestamp(dt.date.today())
    org_df = org_df[org_df['activity_date'] <= as_of].copy()
    if org_df.empty:
        raise ValueError(f"No rows for org '{org_name}' on or before {as_of.date()}")

    # Top-line numbers
    total_events   = int(org_df['event_count'].sum())
    total_projects = int(org_df['project_name'].nunique())
    total_users    = int(org_df['user_name'].nunique())

    # Active projects/users by window
    d30 = as_of - pd.Timedelta(days=30)
    d60 = as_of - pd.Timedelta(days=60)
    recent30 = org_df[org_df['activity_date'] >= d30]
    recent60 = org_df[org_df['activity_date'] >= d60]
    active_projects_30d = int(recent30['project_name'].nunique())
    active_projects_60d = int(recent60['project_name'].nunique())
    active_users_30d    = int(recent30['user_name'].nunique())
    active_users_60d    = int(recent60['user_name'].nunique())
    dormant_projects    = total_projects - active_projects_30d

    # Monthly trend — last 12 months
    org_df['ym'] = org_df['activity_date'].dt.to_period('M')
    monthly = org_df.groupby('ym').agg(
        events=('event_count', 'sum'),
        users=('user_name', 'nunique')
    ).reset_index().sort_values('ym')
    monthly = monthly.tail(12)
    monthly_labels = [str(p) for p in monthly['ym']]
    monthly_events = [int(v) for v in monthly['events']]
    monthly_users  = [int(v) for v in monthly['users']]

    # YoY comparison — prefer last COMPLETE month (exclude current partial month)
    def _month_events(period):
        m = org_df[org_df['ym'] == period]
        return int(m['event_count'].sum())

    all_months = sorted(org_df['ym'].dropna().unique())
    yoy_prev_events = yoy_curr_events = 0
    yoy_prev_label = yoy_curr_label = ""
    yoy_growth_pct = 0.0

    # Current month (partial) excluded from YoY baseline
    current_period = pd.Period(as_of, freq='M')
    complete_months = [m for m in all_months if m < current_period]

    if len(complete_months) >= 13:
        yoy_curr = complete_months[-1]
        yoy_prev = complete_months[-13]
        yoy_curr_events = _month_events(yoy_curr)
        yoy_prev_events = _month_events(yoy_prev)
        if yoy_prev_events > 0:
            yoy_growth_pct = round(100 * (yoy_curr_events - yoy_prev_events) / yoy_prev_events, 1)
        yoy_prev_label = str(yoy_prev)
        yoy_curr_label = str(yoy_curr)

    # Modules
    mod_df = org_df.groupby('module').agg(events=('event_count', 'sum')).reset_index()
    mod_df = mod_df.sort_values('events', ascending=False)
    module_labels = mod_df['module'].astype(str).tolist()
    module_events = [int(v) for v in mod_df['events']]

    # Per-module helpers
    def _mod_sum(name_like):
        mask = org_df['module'].astype(str).str.lower().str.contains(name_like.lower(), na=False)
        return int(org_df[mask]['event_count'].sum())

    def _mod_users(name_like):
        mask = org_df['module'].astype(str).str.lower().str.contains(name_like.lower(), na=False)
        return int(org_df[mask]['user_name'].nunique())

    attendance_total = _mod_sum('attendance')
    attendance_users = _mod_users('attendance')
    task_total       = _mod_sum('task')
    material_total   = _mod_sum('material')
    issue_total      = _mod_sum('issue')

    # Task progress updates (event_type or event_name variants)
    task_rows = org_df[org_df['module'].astype(str).str.lower().str.contains('task', na=False)]
    prog_mask = (
        task_rows['event_type'].astype(str).str.lower().str.contains('progress', na=False) |
        task_rows['event_name'].astype(str).str.lower().str.contains('progress', na=False)
    )
    task_progress_updates = int(task_rows[prog_mask]['event_count'].sum())

    # Top sites
    site_df = org_df.groupby('project_name').agg(
        events=('event_count', 'sum'),
        users=('user_name', 'nunique'),
        last=('activity_date', 'max')
    ).reset_index().sort_values('events', ascending=False)
    top_sites = []
    for _, r in site_df.head(10).iterrows():
        last = r['last']
        days_since = (as_of - last).days if pd.notna(last) else 9999
        top_sites.append({
            'name': str(r['project_name']),
            'events': int(r['events']),
            'users': int(r['users']),
            'days_since': int(days_since),
        })

    # Champions — top users by events (last 60 days only)
    champ_df = recent60.groupby(['user_name', 'role']).agg(
        events=('event_count', 'sum'),
        days=('activity_date', lambda x: x.nunique())
    ).reset_index().sort_values('events', ascending=False)
    champions = []
    for _, r in champ_df.head(6).iterrows():
        champions.append({
            'name': str(r['user_name']),
            'role': str(r['role']),
            'events': int(r['events']),
            'days': int(r['days']),
        })

    # Account meta — ARR, subs_start, renewal_month, KAM, trainer
    def _first_non_null(col):
        s = org_df[col].dropna()
        return s.iloc[0] if not s.empty else None

    arr           = _first_non_null('ARR')
    subs_start    = _first_non_null('subs_start_date')
    renewal_month = _first_non_null('renewal_month')
    kam_email     = _first_non_null('kam_email')
    trainer_email = _first_non_null('trainer_email')

    # Total month span
    months_in_data = len(monthly_events) if monthly_events else 1
    avg_events_per_month = round(sum(monthly_events) / months_in_data, 0) if months_in_data else 0

    data_latest_date = as_of.strftime('%Y-%m-%d') if pd.notna(as_of) else ''

    return {
        'org': org_name,
        'arr': int(arr) if arr is not None and pd.notna(arr) else 0,
        'subs_start': str(subs_start) if subs_start is not None and pd.notna(subs_start) else '',
        'renewal_month': str(renewal_month) if renewal_month is not None and pd.notna(renewal_month) else '',
        'kam_email': str(kam_email) if kam_email else '',
        'trainer_email': str(trainer_email) if trainer_email else '',
        'total_events': total_events,
        'total_projects': total_projects,
        'total_users': total_users,
        'active_projects': active_projects_30d,
        'active_projects_60d': active_projects_60d,
        'active_users_30d': active_users_30d,
        'active_users_60d': active_users_60d,
        'dormant_projects': dormant_projects,
        'monthly_labels': monthly_labels,
        'monthly_events': monthly_events,
        'monthly_users': monthly_users,
        'yoy_prev_label': yoy_prev_label,
        'yoy_curr_label': yoy_curr_label,
        'jan25_events': yoy_prev_events,   # kept for template compatibility
        'jan26_events': yoy_curr_events,
        'yoy_growth_pct': yoy_growth_pct,
        'module_labels': module_labels,
        'module_events': module_events,
        'attendance_total': attendance_total,
        'attendance_users': attendance_users,
        'task_total': task_total,
        'task_progress_updates': task_progress_updates,
        'material_total': material_total,
        'issue_total': issue_total,
        'top_sites': top_sites,
        'champions': champions,
        'avg_events_per_month': avg_events_per_month,
        'data_latest_date': data_latest_date,
    }


# =============================================================================
# CUSTOMIZATION DEFAULTS
# =============================================================================
def _default_customizations(stats):
    """Fallback customizations when the user hasn't supplied any."""
    org_title = stats['org'].title() if stats['org'].isupper() else stats['org']

    return {
        'owner_name': 'Team',
        'owner_title': 'Owner',
        'onboarding_goal': None,              # suppresses goal slide
        'onboarding_goal_short': None,
        'renewal_ask': 'same_arr',            # same_arr | upsell | multi_year | custom
        'renewal_custom_text': None,
        'upsell_amount': None,
        'pain_points': [],                    # empty = skip pain-points slide
        'kam_name': 'Powerplay KAM',
        'kam_email': stats.get('kam_email') or '',
        'trainer_name': 'Powerplay Success Partner',
        'trainer_email': stats.get('trainer_email') or '',
        'company_display_name': org_title,
        'business_description': None,
        'region': None,
        'year1_comparison': None,             # {'events':..,'users':..,'sites':..}
        'review_label': 'Partnership Review',
        'report_date_str': _format_report_date(stats.get('data_latest_date')),
    }


def _format_report_date(iso_date_str):
    if not iso_date_str:
        return dt.date.today().strftime('%B %Y')
    try:
        d = dt.datetime.strptime(iso_date_str, '%Y-%m-%d')
        return d.strftime('%B %Y')
    except Exception:
        return dt.date.today().strftime('%B %Y')


def _merge_customizations(stats, user_overrides):
    base = _default_customizations(stats)
    if user_overrides:
        base.update({k: v for k, v in user_overrides.items() if v is not None})
    return base


# =============================================================================
# SLIDE BUILDERS
# =============================================================================
def _slide_cover(prs, SW, SH, stats, C):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, SW, SH, NAVY_DEEP)
    _add_rect(s, 0, 0, Inches(0.35), SH, AMBER)
    for o in [Inches(0.55), Inches(0.85)]:
        _add_rect(s, SW - Inches(0.25) - o, Inches(-0.5), Inches(0.04), Inches(2.5), AMBER)

    _add_text(s, Inches(0.8), Inches(0.6), Inches(6), Inches(0.4),
              "POWERPLAY  ·  VALUE REVIEW",
              size=12, color=AMBER, bold=True)

    _add_text(s, Inches(0.8), Inches(1.9), Inches(11), Inches(0.4),
              "PREPARED FOR",
              size=11, color=AMBER, bold=True)
    _add_text(s, Inches(0.8), Inches(2.25), Inches(11), Inches(0.55),
              C['owner_name'],
              size=22, bold=True, color=WHITE)

    _add_text(s, Inches(0.8), Inches(3.1), Inches(11), Inches(1.2),
              C['company_display_name'],
              size=52, bold=True, color=WHITE)

    _add_text(s, Inches(0.8), Inches(4.25), Inches(11), Inches(0.6),
              f"{C['review_label']} — {C['report_date_str']}",
              size=22, color=LIGHT_BLUE)

    _add_rect(s, Inches(0.8), Inches(5.05), Inches(1.2), Emu(20000), AMBER)

    subs_note_parts = []
    if stats.get('subs_start'):
        subs_note_parts.append(f"Customer since {stats['subs_start']}")
    if stats.get('renewal_month'):
        subs_note_parts.append(f"Renewal due {stats['renewal_month']}")
    if subs_note_parts:
        _add_text(s, Inches(0.8), Inches(5.25), Inches(11), Inches(0.5),
                  "  ·  ".join(subs_note_parts),
                  size=14, color=LIGHT_BLUE, italic=True)

    presented_by = f"Presented by {C['kam_name']} (KAM)"
    if C['trainer_name']:
        presented_by += f"  &  {C['trainer_name']} (Success Partner)"
    presented_by += "  ·  Powerplay"
    _add_text(s, Inches(0.8), Inches(6.6), Inches(11), Inches(0.35),
              presented_by,
              size=12, color=RGBColor(0x97, 0xA9, 0xDC))


def _slide_exec_summary(prs, SW, SH, stats, C):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, SW, SH, CREAM)
    _add_rect(s, 0, 0, SW, Inches(0.12), AMBER)

    _add_text(s, Inches(0.6), Inches(0.5), Inches(10), Inches(0.5),
              "EXECUTIVE SUMMARY", size=12, color=AMBER, bold=True)

    if C.get('onboarding_goal_short'):
        head = "The goal you set has become a daily site habit."
        sub = (f"You asked for one thing at onboarding — {C['onboarding_goal_short']}. "
               f"Today your team runs {stats['active_projects_60d']} sites on it.")
    else:
        head = f"{stats['total_users']} users. {_fmt_int(stats['total_events'])} activities. One platform."
        sub = (f"Across {stats['total_projects']} project sites, Powerplay is actively used by "
               f"{stats['active_users_30d']} team members this month.")

    _add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
              head, size=32, bold=True, color=NAVY)
    _add_text(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
              sub, size=15, color=MUTED, italic=True)

    # 4 stat cards
    card_y = Inches(2.8)
    card_h = Inches(2.0)
    card_w = Inches(2.95)
    gap    = Inches(0.15)
    start_x = Inches(0.6)

    def stat_card(x, big, big_color, label, sub):
        _add_rect(s, x, card_y, card_w, card_h, WHITE, line=BORDER)
        _add_rect(s, x, card_y, card_w, Inches(0.08), big_color)
        _add_text(s, x + Inches(0.3), card_y + Inches(0.3), card_w - Inches(0.6), Inches(0.8),
                  big, size=42, bold=True, color=big_color)
        _add_text(s, x + Inches(0.3), card_y + Inches(1.05), card_w - Inches(0.6), Inches(0.4),
                  label, size=13, bold=True, color=INK)
        _add_text(s, x + Inches(0.3), card_y + Inches(1.45), card_w - Inches(0.6), Inches(0.5),
                  sub, size=10.5, color=MUTED)

    stat_card(start_x,
              _fmt_int(stats['total_events']), NAVY,
              "Total activities logged",
              f"Across {len(stats['monthly_events'])} months and {stats['total_projects']} sites")
    stat_card(start_x + card_w + gap,
              str(stats['active_projects_60d']), AMBER,
              "Sites active (last 60 days)",
              f"Of {stats['total_projects']} total project sites")
    stat_card(start_x + (card_w + gap) * 2,
              str(stats['active_users_30d']), SUCCESS,
              "Active users this month",
              f"{stats['total_users']} trained users on the account")

    # Fourth card: YoY if available, else total projects
    if stats['yoy_growth_pct'] and stats['yoy_prev_label']:
        stat_card(start_x + (card_w + gap) * 3,
                  f"{'+' if stats['yoy_growth_pct'] >= 0 else ''}{stats['yoy_growth_pct']}%",
                  SUCCESS if stats['yoy_growth_pct'] >= 0 else WARN,
                  "YoY activity growth",
                  f"{stats['yoy_prev_label']} → {stats['yoy_curr_label']}: "
                  f"{_fmt_int(stats['jan25_events'])} → {_fmt_int(stats['jan26_events'])}/mo")
    else:
        stat_card(start_x + (card_w + gap) * 3,
                  str(stats['total_projects']), NAVY,
                  "Projects live on Powerplay",
                  f"{stats['active_projects']} active in the last 30 days")

    # Bottom takeaway bar
    _add_rect(s, Inches(0.6), Inches(5.35), Inches(12.1), Inches(1.55), NAVY)
    _add_rect(s, Inches(0.6), Inches(5.35), Inches(0.1), Inches(1.55), AMBER)
    _add_text(s, Inches(0.95), Inches(5.5), Inches(11.5), Inches(0.4),
              "THE HEADLINE", size=11, bold=True, color=AMBER)

    # Build headline bullets, adapting to customization
    headline = []
    if C.get('onboarding_goal_short'):
        headline.append(f"The original problem you wanted to solve — {C['onboarding_goal_short']} — "
                        f"is fully solved across {stats['active_projects_60d']} active sites.")
    if stats['task_progress_updates'] > 0:
        headline.append(f"Your supervisors are using task progress tracking on their own "
                        f"initiative ({_fmt_int(stats['task_progress_updates'])} updates).")
    headline.append("This is the profile of a successful, sticky deployment — "
                    "the case for renewal writes itself.")

    _add_text(s, Inches(0.95), Inches(5.85), Inches(11.5), Inches(1.0),
              headline, size=14, color=WHITE, spacing_after=2)


def _slide_account_snapshot(prs, SW, SH, stats, C):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, SW, SH, WHITE)
    _add_text(s, Inches(0.6), Inches(0.5), Inches(10), Inches(0.4),
              "01  ·  ACCOUNT SNAPSHOT", size=12, color=AMBER, bold=True)
    _add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
              "Who you are on Powerplay today",
              size=32, bold=True, color=NAVY)

    left_x = Inches(0.6)
    left_w = Inches(6.2)

    rows = [("Company", C['company_display_name'])]
    if C.get('business_description'):
        rows.append(("Business", C['business_description']))
    if C.get('region'):
        rows.append(("Region", C['region']))
    if stats.get('subs_start'):
        rows.append(("Customer since", stats['subs_start']))
    if stats.get('renewal_month'):
        rows.append(("Renewal due", stats['renewal_month']))
    if stats.get('arr'):
        rows.append(("ARR", f"Rs {stats['arr']:,}"))
    rows.append(("Powerplay KAM", C['kam_name']))
    rows.append(("Success partner", C['trainer_name']))

    ry = Inches(2.1)
    for i, (k, v) in enumerate(rows):
        yy = ry + Inches(0.52) * i
        _add_text(s, left_x, yy, Inches(2.0), Inches(0.4), k.upper(),
                  size=10, bold=True, color=MUTED, valign=MSO_ANCHOR.TOP)
        _add_text(s, left_x + Inches(2.1), yy, left_w - Inches(2.1), Inches(0.45), v,
                  size=13, color=INK, valign=MSO_ANCHOR.TOP)

    # Right column — footprint card
    right_x = Inches(7.2)
    right_w = Inches(5.6)
    _add_rect(s, right_x, Inches(2.1), right_w, Inches(4.4), NAVY)
    _add_rect(s, right_x, Inches(2.1), Inches(0.1), Inches(4.4), AMBER)

    _add_text(s, right_x + Inches(0.35), Inches(2.3), right_w - Inches(0.6), Inches(0.45),
              "YOUR POWERPLAY FOOTPRINT", size=11, bold=True, color=AMBER)

    _add_text(s, right_x + Inches(0.35), Inches(2.85), Inches(2.4), Inches(0.8),
              str(stats['total_projects']), size=56, bold=True, color=WHITE)
    _add_text(s, right_x + Inches(0.35), Inches(3.8), Inches(2.4), Inches(0.35),
              "Projects on Powerplay", size=12, color=LIGHT_BLUE)
    _add_text(s, right_x + Inches(0.35), Inches(4.1), Inches(2.4), Inches(0.3),
              f"{stats['active_projects_60d']} currently active",
              size=10, color=AMBER, italic=True)

    _add_text(s, right_x + Inches(2.9), Inches(2.85), Inches(2.4), Inches(0.8),
              str(stats['total_users']), size=56, bold=True, color=WHITE)
    _add_text(s, right_x + Inches(2.9), Inches(3.8), Inches(2.4), Inches(0.35),
              "Trained team members", size=12, color=LIGHT_BLUE)
    _add_text(s, right_x + Inches(2.9), Inches(4.1), Inches(2.4), Inches(0.3),
              f"{stats['active_users_30d']} active in last 30 days",
              size=10, color=AMBER, italic=True)

    _add_text(s, right_x + Inches(0.35), Inches(4.75), right_w - Inches(0.7), Inches(0.8),
              _fmt_int(stats['total_events']), size=50, bold=True, color=WHITE)
    _add_text(s, right_x + Inches(0.35), Inches(5.6), right_w - Inches(0.7), Inches(0.35),
              "Total activities logged across the account",
              size=12, color=LIGHT_BLUE)
    _add_text(s, right_x + Inches(0.35), Inches(5.9), right_w - Inches(0.7), Inches(0.3),
              f"~{int(stats['avg_events_per_month']):,} activities logged every month on average",
              size=10, color=AMBER, italic=True)


def _slide_goal_delivered(prs, SW, SH, stats, C):
    """Optional — only renders if onboarding_goal was supplied."""
    if not C.get('onboarding_goal'):
        return
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, SW, SH, WHITE)
    _add_text(s, Inches(0.6), Inches(0.5), Inches(10), Inches(0.4),
              "02  ·  GOAL DELIVERED", size=12, color=AMBER, bold=True)
    _add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
              "You told us what you wanted. Here's what we delivered.",
              size=30, bold=True, color=NAVY)

    # LEFT — quote on navy
    lx, lw, ly, lh = Inches(0.6), Inches(5.6), Inches(2.3), Inches(4.5)
    _add_rect(s, lx, ly, lw, lh, NAVY)
    _add_rect(s, lx, ly, Inches(0.15), lh, AMBER)

    at_label = "AT ONBOARDING"
    if stats.get('subs_start'):
        at_label = f"AT ONBOARDING — {stats['subs_start']}"
    _add_text(s, lx + Inches(0.5), ly + Inches(0.35), lw - Inches(0.7), Inches(0.4),
              at_label, size=11, bold=True, color=AMBER)
    _add_text(s, lx + Inches(0.5), ly + Inches(0.75), Inches(1), Inches(1),
              "\u201C", size=90, bold=True, color=AMBER)
    _add_text(s, lx + Inches(0.5), ly + Inches(1.9), lw - Inches(0.8), Inches(1.6),
              C['onboarding_goal'],
              size=22, bold=True, color=WHITE, italic=True)
    _add_text(s, lx + Inches(0.5), ly + Inches(3.7), lw - Inches(0.8), Inches(0.45),
              f"— {C['owner_name']}, {C['owner_title']}",
              size=13, color=LIGHT_BLUE, italic=True)

    # RIGHT — metrics
    rx, rw, ry = Inches(6.5), Inches(6.3), Inches(2.3)
    _add_text(s, rx, ry, rw, Inches(0.4),
              "TODAY", size=11, bold=True, color=AMBER)
    _add_text(s, rx, ry + Inches(0.35), rw, Inches(0.6),
              "That single goal is now a company-wide habit.",
              size=18, bold=True, color=NAVY)

    metrics = [
        (_fmt_int(stats['attendance_total']),
         "attendance entries logged across the account"),
        (str(stats['attendance_users']),
         "team members marking attendance on mobile"),
        (str(stats['active_projects_60d']),
         "active project sites running daily"),
        ("0", "paper registers — data flows digitally end-to-end"),
    ]

    my = ry + Inches(1.2)
    for big, small in metrics:
        _add_circle(s, rx + Inches(0.22), my + Inches(0.28), Inches(0.4), SUCCESS)
        _add_text(s, rx + Inches(0.08), my + Inches(0.08), Inches(0.3), Inches(0.4),
                  "\u2713", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_text(s, rx + Inches(0.6), my, Inches(1.8), Inches(0.55),
                  big, size=26, bold=True, color=NAVY)
        _add_text(s, rx + Inches(2.4), my + Inches(0.12), rw - Inches(2.4), Inches(0.45),
                  small, size=13, color=INK)
        my += Inches(0.75)

    _add_rect(s, rx, Inches(6.4), rw, Inches(0.5), AMBER_SOFT)
    _add_rect(s, rx, Inches(6.4), Inches(0.1), Inches(0.5), AMBER)
    _add_text(s, rx + Inches(0.25), Inches(6.5), rw - Inches(0.5), Inches(0.35),
              "Verdict: the problem Powerplay was hired to solve is solved.",
              size=13, bold=True, color=INK)


def _slide_usage_trend(prs, SW, SH, stats, C):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, SW, SH, WHITE)
    _add_text(s, Inches(0.6), Inches(0.5), Inches(10), Inches(0.4),
              "03  ·  USAGE TREND", size=12, color=AMBER, bold=True)

    # Headline depends on YoY trajectory
    if stats['yoy_growth_pct'] and stats['yoy_growth_pct'] > 5:
        head = "Usage has grown every quarter — not plateaued."
    elif stats['yoy_growth_pct'] and stats['yoy_growth_pct'] < -5:
        head = "Usage has declined — worth understanding why."
    else:
        head = "Usage has stabilised — a daily operational habit."
    _add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
              head, size=32, bold=True, color=NAVY)
    _add_text(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
              "Monthly activity logged on Powerplay — last 12 months",
              size=14, color=MUTED, italic=True)

    # Chart
    cd = CategoryChartData()
    months_map = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun',
                  'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
    short_labels = []
    for m in stats['monthly_labels']:
        try:
            y, mo = m.split('-')
            short_labels.append(f"{months_map[int(mo)-1]} {y[-2:]}")
        except Exception:
            short_labels.append(m)
    cd.categories = short_labels
    cd.add_series('Monthly Activities', stats['monthly_events'])

    chart = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.6), Inches(2.55), Inches(8.6), Inches(4.3), cd
    ).chart
    chart.has_title = False
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    dls = plot.data_labels
    dls.font.size = Pt(9)
    dls.font.color.rgb = INK
    dls.position = XL_LABEL_POSITION.OUTSIDE_END

    ser = plot.series[0]
    n = len(stats['monthly_events'])
    for i, pt in enumerate(ser.points):
        fill = pt.format.fill
        fill.solid()
        if i >= n - 4 and i < n - 1:
            fill.fore_color.rgb = AMBER
        elif i == n - 1:
            fill.fore_color.rgb = RGBColor(0xFD, 0xBA, 0x74)
        else:
            fill.fore_color.rgb = NAVY
        pt.format.line.fill.background()

    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(10)
    cat_axis.tick_labels.font.color.rgb = MUTED
    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(10)
    val_axis.tick_labels.font.color.rgb = MUTED
    val_axis.major_gridlines.format.line.color.rgb = BORDER

    # Right insight panel
    px, pw = Inches(9.5), Inches(3.3)
    _add_rect(s, px, Inches(2.55), pw, Inches(4.3), CREAM, line=BORDER)
    _add_rect(s, px, Inches(2.55), pw, Inches(0.08), AMBER)
    _add_text(s, px + Inches(0.3), Inches(2.8), pw - Inches(0.6), Inches(0.4),
              "WHAT THIS SHOWS", size=11, bold=True, color=AMBER)

    yoy_color = SUCCESS if stats['yoy_growth_pct'] >= 0 else WARN
    yoy_label = f"{'+' if stats['yoy_growth_pct'] >= 0 else ''}{stats['yoy_growth_pct']}% YoY"
    _add_text(s, px + Inches(0.3), Inches(3.25), pw - Inches(0.6), Inches(0.5),
              yoy_label, size=24, bold=True, color=yoy_color)
    if stats['yoy_prev_label']:
        _add_text(s, px + Inches(0.3), Inches(3.85), pw - Inches(0.6), Inches(0.6),
                  f"{stats['yoy_prev_label']} → {stats['yoy_curr_label']} monthly "
                  f"activity went from {_fmt_int(stats['jan25_events'])} to "
                  f"{_fmt_int(stats['jan26_events'])}.",
                  size=11, color=INK)

    # Find peak / stable region
    if stats['monthly_events']:
        peak = max(stats['monthly_events'])
        avg  = sum(stats['monthly_events']) / len(stats['monthly_events'])
        _add_text(s, px + Inches(0.3), Inches(4.75), pw - Inches(0.6), Inches(0.45),
                  "Monthly average", size=15, bold=True, color=NAVY)
        _add_text(s, px + Inches(0.3), Inches(5.15), pw - Inches(0.6), Inches(1.0),
                  f"~{_fmt_int(avg)} activities per month. Peak month logged "
                  f"{_fmt_int(peak)} activities.",
                  size=11, color=INK)
    _add_text(s, px + Inches(0.3), Inches(6.3), pw - Inches(0.6), Inches(0.45),
              "Latest month may be in progress",
              size=10, italic=True, color=MUTED)


def _slide_module_breakdown(prs, SW, SH, stats, C):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, SW, SH, WHITE)
    _add_text(s, Inches(0.6), Inches(0.5), Inches(10), Inches(0.4),
              "04  ·  MODULE BREAKDOWN", size=12, color=AMBER, bold=True)

    # Derive headline: top 1-2 modules vs untapped
    mods = list(zip(stats['module_labels'], stats['module_events']))
    mods_filtered = [(l, v) for l, v in mods if v > 10]
    if mods_filtered:
        top = mods_filtered[0][0].title()
        if len(mods_filtered) >= 2:
            second = mods_filtered[1][0].title()
            head = f"{top} & {second} are the backbone. Other modules have room to grow."
        else:
            head = f"{top} is the backbone. Other modules have room to grow."
    else:
        head = "Module usage pattern"
    _add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
              head, size=28, bold=True, color=NAVY)

    # Donut
    cd = CategoryChartData()
    mod_display, mod_values = [], []
    for lbl, v in mods:
        if v > 10:
            mod_display.append(lbl.title())
            mod_values.append(v)
    if not mod_display:
        return
    cd.categories = mod_display
    cd.add_series('Activities by module', mod_values)

    chart = s.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        Inches(0.3), Inches(2.2), Inches(5.5), Inches(4.6), cd
    ).chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(11)
    chart.legend.font.color.rgb = INK
    plot = chart.plots[0]
    plot.has_data_labels = True
    dls = plot.data_labels
    dls.font.size = Pt(11)
    dls.font.color.rgb = WHITE
    dls.font.bold = True
    dls.show_percentage = True
    dls.show_value = False
    dls.show_category_name = False

    palette = [NAVY, AMBER, SUCCESS, MUTED, LIGHT_BLUE]
    ser = plot.series[0]
    for i, pt in enumerate(ser.points):
        fill = pt.format.fill
        fill.solid()
        fill.fore_color.rgb = palette[i % len(palette)]
        pt.format.line.color.rgb = WHITE
        pt.format.line.width = Pt(1.5)

    # Right side — adoption rows
    rx, rw, ry = Inches(6.2), Inches(6.6), Inches(2.25)
    _add_text(s, rx, ry, rw, Inches(0.4), "Adoption snapshot",
              size=16, bold=True, color=NAVY)

    row_defs = [
        ('Attendance', stats['attendance_total'],
         f"{stats['attendance_users']} users have marked attendance — active across sites",
         NAVY),
        ('Task management', stats['task_total'],
         f"{_fmt_int(stats['task_progress_updates'])} progress updates logged by supervisors"
         if stats['task_progress_updates'] else
         "Task module in use — progress tracking available",
         AMBER),
        ('Material management', stats['material_total'],
         "GRN, indents, transfers — module is in active use"
         if stats['material_total'] > max(500, stats['attendance_total'] * 0.3) else
         "Limited usage. GRN, indents, transfers largely unused — expansion opportunity",
         SUCCESS),
        ('Issue tracker', stats['issue_total'],
         "Active — structured issue/RFI tracking in place"
         if stats['issue_total'] > 200 else
         "Nearly dormant — no structured issue/RFI tracking happening today",
         MUTED),
    ]

    rry = Inches(2.8)
    for name, total, desc, clr in row_defs:
        _add_rect(s, rx, rry, Inches(0.1), Inches(0.9), clr)
        _add_text(s, rx + Inches(0.3), rry, Inches(3.5), Inches(0.4),
                  name, size=15, bold=True, color=INK)
        _add_text(s, rx + rw - Inches(1.8), rry, Inches(1.8), Inches(0.4),
                  f"{_fmt_int(total)} events", size=15, bold=True, color=clr,
                  align=PP_ALIGN.RIGHT)
        _add_text(s, rx + Inches(0.3), rry + Inches(0.42), rw - Inches(0.4), Inches(0.5),
                  desc, size=11, color=MUTED)
        rry += Inches(1.05)


def _slide_top_sites(prs, SW, SH, stats, C):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, SW, SH, WHITE)
    _add_text(s, Inches(0.6), Inches(0.5), Inches(10), Inches(0.4),
              "05  ·  SITE-LEVEL ADOPTION", size=12, color=AMBER, bold=True)

    active = stats['active_projects_60d']
    head = f"{active} project sites run on Powerplay regularly."
    _add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
              head, size=30, bold=True, color=NAVY)
    _add_text(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
              f"Top {min(len(stats['top_sites']), 10)} sites by total activity "
              f"— ordered by depth of use",
              size=14, color=MUTED, italic=True)

    # Bar chart
    cd = CategoryChartData()
    sites = list(reversed(stats['top_sites'][:10]))
    cd.categories = [s_['name'].replace('_', ' ') for s_ in sites]
    cd.add_series('Events', [s_['events'] for s_ in sites])

    chart = s.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.6), Inches(2.5), Inches(8.8), Inches(4.6), cd
    ).chart
    chart.has_title = False
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    dls = plot.data_labels
    dls.font.size = Pt(10)
    dls.font.color.rgb = INK
    dls.position = XL_LABEL_POSITION.OUTSIDE_END
    ser = plot.series[0]
    for i, pt in enumerate(ser.points):
        fill = pt.format.fill
        fill.solid()
        fill.fore_color.rgb = NAVY
        pt.format.line.fill.background()
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(10)
    cat_axis.tick_labels.font.color.rgb = INK
    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(9)
    val_axis.tick_labels.font.color.rgb = MUTED
    val_axis.major_gridlines.format.line.color.rgb = BORDER

    # Right panel
    px, pw = Inches(9.7), Inches(3.2)
    _add_rect(s, px, Inches(2.5), pw, Inches(4.6), CREAM, line=BORDER)
    _add_rect(s, px, Inches(2.5), pw, Inches(0.08), AMBER)
    _add_text(s, px + Inches(0.25), Inches(2.75), pw - Inches(0.5), Inches(0.4),
              "SITE HEALTH", size=11, bold=True, color=AMBER)
    _add_text(s, px + Inches(0.25), Inches(3.2), pw - Inches(0.5), Inches(0.5),
              f"{stats['active_projects']} / {stats['total_projects']}",
              size=30, bold=True, color=SUCCESS)
    _add_text(s, px + Inches(0.25), Inches(3.85), pw - Inches(0.5), Inches(0.4),
              "sites active in last 30 days", size=11, color=INK)

    _add_text(s, px + Inches(0.25), Inches(4.4), pw - Inches(0.5), Inches(0.4),
              "Adoption leaders", size=13, bold=True, color=NAVY)

    top3 = stats['top_sites'][:3] if stats['top_sites'] else []
    leaders_line = ", ".join([t['name'].replace('_', ' ')[:30] for t in top3])
    body_lines = [
        f"{leaders_line} lead adoption with the most activities." if leaders_line else "",
        "",
        f"{stats['dormant_projects']} sites are dormant — natural if they are "
        f"closed/completed, but worth reviewing."
        if stats['dormant_projects'] > 0 else "",
    ]
    _add_text(s, px + Inches(0.25), Inches(4.75), pw - Inches(0.5), Inches(2.0),
              [l for l in body_lines if l is not None],
              size=10, color=INK, spacing_after=3)


def _slide_champions(prs, SW, SH, stats, C):
    if not stats['champions']:
        return
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, SW, SH, CREAM)
    _add_text(s, Inches(0.6), Inches(0.5), Inches(10), Inches(0.4),
              "06  ·  YOUR POWERPLAY CHAMPIONS", size=12, color=AMBER, bold=True)
    _add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
              "The supervisors who made Powerplay a daily habit.",
              size=30, bold=True, color=NAVY)
    _add_text(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
              "Top active users in the last 60 days — the people driving value on the ground.",
              size=14, color=MUTED, italic=True)

    card_w = Inches(4.05)
    card_h = Inches(2.1)
    gap    = Inches(0.15)
    start_x, start_y = Inches(0.6), Inches(2.65)

    for i, c in enumerate(stats['champions'][:6]):
        col, row = i % 3, i // 3
        x = start_x + (card_w + gap) * col
        y = start_y + (card_h + gap) * row
        _add_rect(s, x, y, card_w, card_h, WHITE, line=BORDER)
        _add_rect(s, x, y, card_w, Inches(0.08), AMBER)

        init = ''.join([w[0] for w in c['name'].split()[:2]]).upper() or '—'
        _add_circle(s, x + Inches(0.7), y + Inches(0.85), Inches(0.9), NAVY)
        _add_text(s, x + Inches(0.25), y + Inches(0.65), Inches(0.9), Inches(0.5),
                  init, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        _add_text(s, x + Inches(1.3), y + Inches(0.4), card_w - Inches(1.5), Inches(0.4),
                  c['name'], size=16, bold=True, color=INK)
        role_clean = c['role'].split('_')[0] if '_' in c['role'] else c['role']
        _add_text(s, x + Inches(1.3), y + Inches(0.78), card_w - Inches(1.5), Inches(0.35),
                  role_clean, size=11, color=MUTED, italic=True)

        _add_text(s, x + Inches(0.3), y + Inches(1.45), Inches(1.6), Inches(0.4),
                  _fmt_int(c['events']), size=20, bold=True, color=NAVY)
        _add_text(s, x + Inches(0.3), y + Inches(1.75), Inches(1.6), Inches(0.3),
                  "activities", size=10, color=MUTED)
        _add_text(s, x + Inches(2.2), y + Inches(1.45), Inches(1.6), Inches(0.4),
                  str(c['days']), size=20, bold=True, color=AMBER)
        _add_text(s, x + Inches(2.2), y + Inches(1.75), Inches(1.6), Inches(0.3),
                  "active days", size=10, color=MUTED)


def _slide_value_realised(prs, SW, SH, stats, C):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, SW, SH, WHITE)
    _add_text(s, Inches(0.6), Inches(0.5), Inches(10), Inches(0.4),
              "07  ·  VALUE REALISED", size=12, color=AMBER, bold=True)
    _add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
              "What Powerplay is doing for your business right now.",
              size=30, bold=True, color=NAVY)

    pillars = []
    if stats['attendance_total'] > 100:
        pillars.append({
            "icon_color": NAVY,
            "title": f"Zero-paper attendance across {stats['active_projects_60d']} sites",
            "stat": f"{_fmt_int(stats['attendance_total'])} attendance entries",
            "body": "Site supervisors mark attendance directly on mobile — no muster "
                    "registers, no lag in payroll inputs, no disputes over overtime. You "
                    "can see exactly who was on which site, on which day."
        })
    if stats['task_total'] > 100:
        pillars.append({
            "icon_color": AMBER,
            "title": "Live visibility into site progress",
            "stat": f"{_fmt_int(stats['task_progress_updates'] or stats['task_total'])} progress updates",
            "body": "Every update is timestamped and tied to a site and a task. "
                    "Senior management can review what is moving and what is stuck "
                    "without chasing supervisors over phone calls."
        })
    pillars.append({
        "icon_color": SUCCESS,
        "title": "Consistent working habit",
        "stat": f"{stats['active_users_60d']} active users, 60 days",
        "body": "A stable pool of users works on Powerplay every week. The tool is "
                "embedded in how your sites run — not a one-time rollout that faded."
    })
    pillars = pillars[:3]

    py, ph = Inches(2.3), Inches(4.5)
    pw, gx = Inches(4.05), Inches(0.15)
    sx = Inches(0.6)

    for i, p in enumerate(pillars):
        x = sx + (pw + gx) * i
        _add_rect(s, x, py, pw, ph, CREAM, line=BORDER)
        _add_rect(s, x, py, pw, Inches(0.1), p['icon_color'])
        _add_text(s, x + Inches(0.35), py + Inches(0.35), pw - Inches(0.7), Inches(0.5),
                  f"0{i+1}", size=14, bold=True, color=p['icon_color'])
        _add_text(s, x + Inches(0.35), py + Inches(0.75), pw - Inches(0.7), Inches(1.1),
                  p['title'], size=18, bold=True, color=INK)
        _add_text(s, x + Inches(0.35), py + Inches(2.15), pw - Inches(0.7), Inches(0.6),
                  p['stat'], size=22, bold=True, color=p['icon_color'])
        _add_text(s, x + Inches(0.35), py + Inches(2.85), pw - Inches(0.7), Inches(1.5),
                  p['body'], size=11, color=INK)


def _slide_pain_points(prs, SW, SH, stats, C):
    """Pain-points slide — only renders when pain_points list is non-empty."""
    pps = C.get('pain_points') or []
    if not pps:
        return
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, SW, SH, WHITE)
    _add_text(s, Inches(0.6), Inches(0.5), Inches(10), Inches(0.4),
              "08  ·  WHAT YOU'VE FLAGGED", size=12, color=AMBER, bold=True)
    _add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
              "What you've told us — and how Year 3 will address it.",
              size=28, bold=True, color=NAVY)
    _add_text(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
              "Every pain point you raised becomes a concrete commitment in our Year 3 plan.",
              size=14, color=MUTED, italic=True)

    ry = Inches(2.55)
    rh = Inches(0.95)
    for i, p in enumerate(pps[:5]):
        y = ry + (rh + Inches(0.12)) * i
        _add_rect(s, Inches(0.6), y, Inches(12.1), rh, CREAM, line=BORDER)
        _add_rect(s, Inches(0.6), y, Inches(0.15), rh, AMBER)
        # Warning icon
        _add_circle(s, Inches(1.0), y + Inches(0.47), Inches(0.45), WARN)
        _add_text(s, Inches(0.85), y + Inches(0.29), Inches(0.3), Inches(0.4),
                  "!", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_text(s, Inches(1.55), y + Inches(0.15), Inches(3.0), Inches(0.35),
                  "WHAT YOU FLAGGED", size=9, bold=True, color=MUTED)
        _add_text(s, Inches(1.55), y + Inches(0.45), Inches(4.5), Inches(0.5),
                  str(p), size=14, bold=True, color=INK)
        _add_text(s, Inches(6.2), y + Inches(0.15), Inches(3.0), Inches(0.35),
                  "HOW WE'LL ADDRESS IT", size=9, bold=True, color=AMBER)
        _add_text(s, Inches(6.2), y + Inches(0.45), Inches(6.3), Inches(0.5),
                  "Covered under the Year 3 commitment (next slide) — "
                  "specific actions in May–July 2026.",
                  size=12, color=INK)


def _slide_whats_next(prs, SW, SH, stats, C):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, SW, SH, WHITE)
    _add_text(s, Inches(0.6), Inches(0.5), Inches(10), Inches(0.4),
              "09  ·  WHAT'S NEXT FOR YEAR 3", size=12, color=AMBER, bold=True)
    _add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
              "Three areas where we can deliver more value next year.",
              size=28, bold=True, color=NAVY)

    opps = []
    if stats['material_total'] < stats['attendance_total'] * 0.5:
        opps.append({
            "title": "Material & inventory management",
            "state": f"Currently under-used — ~{_fmt_int(stats['material_total'])} events "
                     f"vs {_fmt_int(stats['attendance_total'])} on attendance",
            "impact": "Wastage visibility & payment hygiene",
            "detail": "Activate indents, GRN, site-to-site transfers and MIS reports. "
                      "Even 1–2% reduction in material leakage on a multi-project "
                      "portfolio pays back the subscription many times over.",
            "color": AMBER
        })
    if stats['issue_total'] < 100:
        opps.append({
            "title": "Issue / RFI tracking",
            "state": f"Currently dormant — only {stats['issue_total']} events logged",
            "impact": "Faster decisions, fewer site-to-HO phone calls",
            "detail": "Structured issue logs create an audit trail for decisions taken "
                      "on site. Useful for govt/PSU projects where RFIs need documentation.",
            "color": NAVY
        })
    if stats['dormant_projects'] > 0:
        opps.append({
            "title": "Bring dormant sites back online",
            "state": f"{stats['dormant_projects']} sites with no activity in 30+ days",
            "impact": "Full-portfolio visibility for leadership",
            "detail": "Some may be closed/complete — but any live site not on "
                      "Powerplay is a blind spot. Worth a 30-min review with our team.",
            "color": SUCCESS
        })

    # Fallback: always have at least 3
    while len(opps) < 3:
        opps.append({
            "title": "Advanced analytics & reports",
            "state": "Standard reports in place",
            "impact": "Deeper leadership visibility",
            "detail": "Custom dashboards, labour productivity metrics, project P&L "
                      "views — available in higher tiers.",
            "color": AMBER
        })
    opps = opps[:3]

    py = Inches(2.2)
    ph = Inches(1.55)
    for i, o in enumerate(opps):
        y = py + (ph + Inches(0.12)) * i
        _add_rect(s, Inches(0.6), y, Inches(12.1), ph, CREAM, line=BORDER)
        _add_rect(s, Inches(0.6), y, Inches(0.15), ph, o['color'])
        _add_text(s, Inches(0.95), y + Inches(0.2), Inches(4.8), Inches(0.5),
                  o['title'], size=16, bold=True, color=INK)
        _add_text(s, Inches(0.95), y + Inches(0.7), Inches(4.8), Inches(0.4),
                  o['state'], size=11, color=MUTED, italic=True)
        _add_text(s, Inches(5.9), y + Inches(0.2), Inches(2.5), Inches(0.35),
                  "BUSINESS IMPACT", size=9, bold=True, color=o['color'])
        _add_text(s, Inches(5.9), y + Inches(0.5), Inches(2.5), Inches(0.9),
                  o['impact'], size=13, bold=True, color=INK)
        _add_text(s, Inches(8.6), y + Inches(0.2), Inches(4.0), Inches(0.35),
                  "HOW TO ACTIVATE", size=9, bold=True, color=o['color'])
        _add_text(s, Inches(8.6), y + Inches(0.5), Inches(4.0), Inches(1.1),
                  o['detail'], size=10, color=INK)


def _slide_year1_comparison(prs, SW, SH, stats, C):
    """Y1 vs Y2 comparison — only renders when year1_comparison is supplied."""
    y1 = C.get('year1_comparison')
    if not y1 or not isinstance(y1, dict):
        return

    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, SW, SH, CREAM)
    _add_text(s, Inches(0.6), Inches(0.5), Inches(10), Inches(0.4),
              "10  ·  YEAR 1 vs YEAR 2", size=12, color=AMBER, bold=True)
    _add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
              "How your adoption curve has moved year over year.",
              size=30, bold=True, color=NAVY)

    metrics = []
    if y1.get('events') is not None:
        metrics.append(('Total activities', int(y1['events']), stats['total_events']))
    if y1.get('users') is not None:
        metrics.append(('Trained users', int(y1['users']), stats['total_users']))
    if y1.get('sites') is not None:
        metrics.append(('Active sites', int(y1['sites']), stats['active_projects_60d']))

    if not metrics:
        return

    card_w = Inches(4.05)
    card_h = Inches(3.9)
    gap    = Inches(0.15)
    start_x = Inches(0.6)
    card_y = Inches(2.5)

    for i, (label, y1v, y2v) in enumerate(metrics[:3]):
        x = start_x + (card_w + gap) * i
        _add_rect(s, x, card_y, card_w, card_h, WHITE, line=BORDER)
        _add_rect(s, x, card_y, card_w, Inches(0.1), AMBER)

        _add_text(s, x + Inches(0.35), card_y + Inches(0.35), card_w - Inches(0.7),
                  Inches(0.4), label.upper(), size=11, bold=True, color=AMBER)

        # Y1
        _add_text(s, x + Inches(0.35), card_y + Inches(0.85), card_w - Inches(0.7),
                  Inches(0.35), "YEAR 1", size=10, bold=True, color=MUTED)
        _add_text(s, x + Inches(0.35), card_y + Inches(1.15), card_w - Inches(0.7),
                  Inches(0.7), _fmt_int(y1v), size=28, bold=True, color=MUTED)

        # Y2
        _add_text(s, x + Inches(0.35), card_y + Inches(2.05), card_w - Inches(0.7),
                  Inches(0.35), "YEAR 2", size=10, bold=True, color=NAVY)
        _add_text(s, x + Inches(0.35), card_y + Inches(2.35), card_w - Inches(0.7),
                  Inches(0.7), _fmt_int(y2v), size=34, bold=True, color=NAVY)

        # Change pill
        delta = y2v - y1v
        pct = round(100 * delta / y1v, 1) if y1v else 0
        pill_color = SUCCESS if delta >= 0 else WARN
        _add_rect(s, x + Inches(0.35), card_y + Inches(3.15),
                  card_w - Inches(0.7), Inches(0.45), pill_color)
        sign = '+' if delta >= 0 else ''
        _add_text(s, x + Inches(0.35), card_y + Inches(3.23),
                  card_w - Inches(0.7), Inches(0.3),
                  f"{sign}{_fmt_int(delta)}  ·  {sign}{pct}%",
                  size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def _slide_commitment(prs, SW, SH, stats, C):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, SW, SH, CREAM)
    _add_text(s, Inches(0.6), Inches(0.5), Inches(10), Inches(0.4),
              "11  ·  OUR COMMITMENT", size=12, color=AMBER, bold=True)
    _add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
              "Once you renew, here's the first 90 days — on us.",
              size=30, bold=True, color=NAVY)

    kam = C['kam_name']
    trainer = C['trainer_name']
    intro = (f"Commitments from {kam}"
             + (f" and {trainer}" if trainer else "")
             + " — tracked jointly with your team.")
    _add_text(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
              intro, size=14, color=MUTED, italic=True)

    # Build phase data from customizations / stats
    month_labels = _get_next_3_months()
    phases = [
        {
            "phase": month_labels[0],
            "headline": "Material module activation"
                        if stats['material_total'] < stats['attendance_total'] * 0.5
                        else "Module deep-dive",
            "actions": [
                "On-site training with storekeepers at HO + top 3 sites",
                "GRN and indent workflow rollout for lead sites",
                "Weekly report template shared with PM team",
            ],
            "color": AMBER
        },
        {
            "phase": month_labels[1],
            "headline": "Issue tracking + dormant sites",
            "actions": [
                "Light-touch Issue/RFI workflow launched for 5 pilot sites",
                "Review of dormant projects with PM — reactivate or archive",
                "Mid-quarter usage checkpoint with KAM",
            ],
            "color": NAVY
        },
        {
            "phase": month_labels[2],
            "headline": "Scale & measure",
            "actions": [
                "Extend Material + Issue adoption to all active sites",
                "Second Value Review with numbers vs today's baseline",
                "Explore project P&L / payments module pilot",
            ],
            "color": SUCCESS
        },
    ]

    py = Inches(2.7)
    ph = Inches(4.3)
    pw, gx = Inches(4.05), Inches(0.15)
    sx = Inches(0.6)

    for i, p in enumerate(phases):
        x = sx + (pw + gx) * i
        _add_rect(s, x, py, pw, ph, WHITE, line=BORDER)
        _add_rect(s, x, py, pw, Inches(0.5), p['color'])
        _add_text(s, x + Inches(0.35), py + Inches(0.1), pw - Inches(0.7), Inches(0.35),
                  p['phase'], size=12, bold=True, color=WHITE)
        _add_text(s, x + Inches(0.35), py + Inches(0.75), pw - Inches(0.7), Inches(0.9),
                  p['headline'], size=17, bold=True, color=INK)
        by = py + Inches(1.8)
        for act in p['actions']:
            _add_circle(s, x + Inches(0.45), by + Inches(0.14), Inches(0.12), p['color'])
            _add_text(s, x + Inches(0.7), by, pw - Inches(0.9), Inches(0.8),
                      act, size=11, color=INK)
            by += Inches(0.75)


def _get_next_3_months():
    """Returns next 3 month labels based on today, e.g. ['MAY 2026', 'JUNE 2026', 'JULY 2026']."""
    today = dt.date.today()
    labels = []
    m = today.month
    y = today.year
    for i in range(3):
        nm = ((m - 1 + i + 1) % 12) + 1   # next month and 2 after
        ny = y + ((m - 1 + i + 1) // 12)
        labels.append(dt.date(ny, nm, 1).strftime('%B %Y').upper())
    return labels


def _slide_ask(prs, SW, SH, stats, C):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, SW, SH, NAVY_DEEP)
    _add_rect(s, 0, 0, Inches(0.35), SH, AMBER)
    for o in [Inches(0.55), Inches(0.85), Inches(1.15)]:
        _add_rect(s, SW - o, Inches(-0.5), Inches(0.04), Inches(2.5), AMBER)

    _add_text(s, Inches(0.9), Inches(0.7), Inches(11), Inches(0.5),
              "12  ·  OUR ASK", size=12, bold=True, color=AMBER)

    name = C['owner_name'].split()[0] if C['owner_name'] else 'Team'
    _add_text(s, Inches(0.9), Inches(1.25), Inches(11), Inches(0.6),
              f"{name}, we'd like to continue the partnership.",
              size=30, bold=True, color=WHITE)

    # Ask card — varies by renewal_ask
    ask = C.get('renewal_ask', 'same_arr')
    arr = stats.get('arr') or 0
    if ask == 'same_arr':
        ask_headline = f"Renew Powerplay for Year 3 at the same annual fee — Rs {arr:,}."
        ask_sub = "No price change. Same team, same support, same platform."
    elif ask == 'upsell':
        amt = C.get('upsell_amount') or arr
        ask_headline = f"Renew + upgrade for Year 3 at Rs {int(amt):,}."
        ask_sub = "Expanded modules, priority support, and the roadmap on slide 11."
    elif ask == 'multi_year':
        ask_headline = f"Lock in a 2-year renewal at Rs {arr:,}/year."
        ask_sub = "Price protected for 24 months. Same team, extended commitment."
    elif ask == 'custom' and C.get('renewal_custom_text'):
        ask_headline = C['renewal_custom_text']
        ask_sub = ""
    else:
        ask_headline = f"Renew Powerplay for Year 3 at Rs {arr:,}."
        ask_sub = "Continuation of the current plan."

    _add_rect(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.8), AMBER)
    _add_text(s, Inches(1.2), Inches(2.5), Inches(11), Inches(0.4),
              "THE RENEWAL", size=11, bold=True, color=NAVY_DEEP)
    _add_text(s, Inches(1.2), Inches(2.85), Inches(11), Inches(0.65),
              ask_headline, size=22, bold=True, color=NAVY_DEEP)
    if ask_sub:
        _add_text(s, Inches(1.2), Inches(3.55), Inches(11), Inches(0.4),
                  ask_sub, size=13, color=NAVY_DEEP, italic=True)

    # 3 reasons
    rx = Inches(0.9)
    ry = Inches(4.5)
    rw = Inches(3.77)
    gap = Inches(0.15)

    r1 = ("Your goal is being met",
          f"{_fmt_int(stats['attendance_total'])} attendance entries, "
          f"{stats['active_projects_60d']} live sites.") \
         if C.get('onboarding_goal_short') else \
         ("Strong, sticky adoption",
          f"{stats['active_users_30d']} monthly active users on {stats['total_projects']} sites.")

    if stats['yoy_growth_pct'] > 0:
        r2 = ("Adoption is still growing",
              f"+{stats['yoy_growth_pct']}% YoY. Your team is using it more, not less.")
    else:
        r2 = ("Mature, stable usage",
              f"{_fmt_int(stats['total_events'])} total activities — platform is embedded.")
    r3 = ("Year 3 is value, not cost",
          "Material, Issue, dormant sites — more ROI coming.")

    reasons = [r1, r2, r3]
    for i, (t, d) in enumerate(reasons):
        x = rx + (rw + gap) * i
        _add_rect(s, x, ry, rw, Inches(1.2), NAVY)
        _add_rect(s, x, ry, Inches(0.08), Inches(1.2), AMBER)
        _add_text(s, x + Inches(0.25), ry + Inches(0.15), rw - Inches(0.4), Inches(0.4),
                  t, size=14, bold=True, color=WHITE)
        _add_text(s, x + Inches(0.25), ry + Inches(0.55), rw - Inches(0.4), Inches(0.6),
                  d, size=11, color=LIGHT_BLUE)

    # Bottom contact strip
    _add_text(s, Inches(0.9), Inches(6.2), Inches(11), Inches(0.35),
              "YOUR POWERPLAY TEAM — HERE FOR ANYTHING YOU NEED",
              size=10, bold=True, color=AMBER)
    _add_text(s, Inches(0.9), Inches(6.55), Inches(6), Inches(0.3),
              f"{C['kam_name']}  ·  Key Account Manager  ·  {C['kam_email']}",
              size=12, color=WHITE)
    if C['trainer_name']:
        _add_text(s, Inches(0.9), Inches(6.85), Inches(6), Inches(0.3),
                  f"{C['trainer_name']}  ·  Customer Success Partner  ·  {C['trainer_email']}",
                  size=12, color=WHITE)

    _add_text(s, Inches(7.5), Inches(6.55), Inches(5.4), Inches(0.7),
              [f"Thank you, {name}, for your trust.",
               "We're excited for Year 3 together."],
              size=12, color=LIGHT_BLUE, italic=True,
              align=PP_ALIGN.RIGHT, spacing_after=2)


# =============================================================================
# BUILD ONE DECK — orchestrates all slides
# =============================================================================
def build_deck(stats, customizations, output_path):
    C = _merge_customizations(stats, customizations)

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height

    # Slide order. Optional slides silently skip when inputs are missing.
    _slide_cover(prs, SW, SH, stats, C)
    _slide_exec_summary(prs, SW, SH, stats, C)
    _slide_account_snapshot(prs, SW, SH, stats, C)
    _slide_goal_delivered(prs, SW, SH, stats, C)          # optional
    _slide_usage_trend(prs, SW, SH, stats, C)
    _slide_module_breakdown(prs, SW, SH, stats, C)
    _slide_top_sites(prs, SW, SH, stats, C)
    _slide_champions(prs, SW, SH, stats, C)
    _slide_value_realised(prs, SW, SH, stats, C)
    _slide_pain_points(prs, SW, SH, stats, C)             # optional
    _slide_whats_next(prs, SW, SH, stats, C)
    _slide_year1_comparison(prs, SW, SH, stats, C)        # optional
    _slide_commitment(prs, SW, SH, stats, C)
    _slide_ask(prs, SW, SH, stats, C)

    prs.save(output_path)
    return output_path, len(prs.slides)


# =============================================================================
# PUBLIC API — generate_vr_deck
# =============================================================================
def list_orgs(csv_path):
    """Return the list of unique orgs in a CSV."""
    df = pd.read_csv(csv_path, low_memory=False)
    if 'org_name' not in df.columns:
        raise ValueError("CSV is missing required column: org_name")
    return sorted(df['org_name'].dropna().unique().tolist())


def generate_vr_deck(csv_path, org_filter=None, output_dir='./vr_output',
                      customizations=None, save_stats=False, verbose=True):
    """
    Generate VR decks from a usage CSV.

    Args:
        csv_path (str): Path to the depth-funnel CSV.
        org_filter (str | list[str] | None): If None, generate for every org in
            the CSV. If a string, generate for that one org. If a list, that
            subset of orgs.
        output_dir (str): Folder where the .pptx files go (created if missing).
        customizations (dict): Mapping of org_name → per-org customization dict.
            See module docstring for the full key set. Keys are optional; any
            unsupplied key gets a sensible default.
        save_stats (bool): If True, also writes the computed stats dict as JSON
            next to the .pptx — useful for debugging.
        verbose (bool): Print progress to stdout.

    Returns:
        list[dict]: one entry per deck, each with 'org', 'path', 'slides'.
    """
    os.makedirs(output_dir, exist_ok=True)
    customizations = customizations or {}

    if verbose:
        print(f"Reading {csv_path} ...")
    df = pd.read_csv(csv_path, low_memory=False)

    required = {'org_name', 'project_name', 'user_name', 'role',
                'module', 'event_count', 'activity_date'}
    missing = required - set(df.columns)
    if missing:
        raise ValueError(f"CSV is missing columns: {sorted(missing)}")

    # Decide which orgs to process
    all_orgs = sorted(df['org_name'].dropna().unique().tolist())
    if org_filter is None:
        orgs = all_orgs
    elif isinstance(org_filter, str):
        orgs = [org_filter]
    else:
        orgs = list(org_filter)

    # Validate
    missing_orgs = [o for o in orgs if o not in all_orgs]
    if missing_orgs:
        raise ValueError(f"Orgs not found in CSV: {missing_orgs}\n"
                         f"Available: {all_orgs[:10]}{'...' if len(all_orgs) > 10 else ''}")

    results = []
    for org in orgs:
        if verbose:
            print(f"  → {org}")
        stats = compile_stats(df, org)
        out_name = _safe_filename(org)
        date_str = dt.date.today().strftime('%Y_%m')
        out_path = os.path.join(output_dir, f"{out_name}_VR_{date_str}.pptx")
        path, n_slides = build_deck(stats, customizations.get(org), out_path)
        results.append({'org': org, 'path': path, 'slides': n_slides})
        if save_stats:
            stats_path = os.path.join(output_dir, f"{out_name}_stats.json")
            with open(stats_path, 'w') as f:
                json.dump(stats, f, indent=2, default=str)
        if verbose:
            print(f"     saved → {path}  ({n_slides} slides)")

    return results


def zip_decks(paths, zip_path):
    """Convenience: bundle a list of .pptx paths into a zip."""
    with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as z:
        for p in paths:
            z.write(p, arcname=os.path.basename(p))
    return zip_path


# =============================================================================
# CLI
# =============================================================================
if __name__ == '__main__':
    import argparse
    ap = argparse.ArgumentParser(description='Generate Powerplay VR decks from a usage CSV.')
    ap.add_argument('csv', help='Path to usage CSV')
    ap.add_argument('--org', help='Single org name (omit to process every org)')
    ap.add_argument('--out', default='./vr_output', help='Output folder')
    ap.add_argument('--customizations', help='Path to customizations JSON')
    ap.add_argument('--save-stats', action='store_true', help='Also dump stats JSON')
    ap.add_argument('--list-orgs', action='store_true', help='Just list orgs and exit')
    args = ap.parse_args()

    if args.list_orgs:
        for o in list_orgs(args.csv):
            print(o)
    else:
        customs = None
        if args.customizations:
            with open(args.customizations) as f:
                customs = json.load(f)
        generate_vr_deck(args.csv, org_filter=args.org, output_dir=args.out,
                         customizations=customs, save_stats=args.save_stats)
